from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def make_presentation(presentation_title, slide_contents):

    # Create a presentation object
    prs = Presentation()

    ############################################################################################################### TITLE SLIDE #####

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

    title.text = presentation_title.upper()
    subtitle.text = "The Last Bear Ender Artbook"

    # Change title font color to white
    for run in title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = Pt(56)  # Set the font size

    # Change subtitle font color to white
    for run in subtitle.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(100, 100, 100)
        run.font.size = Pt(20)  # Set the font size

    # Add the image
    image_left = Inches(2)
    image_top = Inches(1)
    image_height = Inches(2.5)  # Width will be scaled proportionally
    title_slide.shapes.add_picture('logo_white.png', image_left, image_top, height=image_height)

    ############################################################################################################### EMPTY SLIDE #####

    second_slide_layout = prs.slide_layouts[6]
    black_slide = prs.slides.add_slide(second_slide_layout)

    ############################################################################################################### NEXT STEP IN PROCESS SLIDE #####

    for artstage in slide_contents:
        # Add a black slide
        black_slide_layout = prs.slide_layouts[6]  # Using a blank layout
        black_slide = prs.slides.add_slide(black_slide_layout)
        background = black_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

        # Add the text box
        text_box = black_slide.shapes.add_textbox(Inches(1), Inches(6), Inches(6), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = artstage["title"]

        # Format the text
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(44)  # Set the font size
                run.font.color.rgb = RGBColor(255, 255, 255)  # White color
                run.font.bold = True

        for image_url in artstage["images"]:
            ############################################################################################################### NEW IMAGE SLIDE #####

            # Add a slide without a title but with an image and text
            image_slide_layout = prs.slide_layouts[6]  # Using a blank layout
            image_slide = prs.slides.add_slide(image_slide_layout)

            # Add the image
            image_left = Inches(2)
            image_top = Inches(1)
            image_height = Inches(2.5)  # Width will be scaled proportionally
            print("IMAGE: " + image_url)
            image_slide.shapes.add_picture(image_url, image_left, image_top, height=image_height)

    ############################################################################################################### CLOSING SLIDE #####

    # Add a closing slide with an image
    closing_slide_layout = prs.slide_layouts[6]  # Choosing a blank layout for the closing slide
    closing_slide = prs.slides.add_slide(closing_slide_layout)
    logo_path = 'logo_black.png'  # Replace with the path to your logo image

    # Add the logo image to the slide
    logo_left = Inches(2)  # Adjust the position as needed
    logo_top = Inches(1)  # Adjust the position as needed
    logo_height = Inches(2.5)  # The width will be scaled proportionally
    closing_slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)

    # Add the disclaimer text at the bottom
    disclaimer_text = ("Disclaimer: All images presented herein were created using DALL·E technology."
                    "Wilrik De Loose, also known as The Last Bear Ender, holds and retains all copyright ownership of these images."
                    "Unauthorized use, reproduction, or generation of images based on the works of The Last Bear Ender, including the images generated and displayed here, is strictly prohibited."
                    "All rights are reserved by Wilrik De Loose, and any infringement of these rights will be fully pursued of the law.")

    text_box = closing_slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = disclaimer_text

    # Format the disclaimer text
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(128, 128, 128)  # Grey color

    # Save the presentation
    prs.save('Beareaved.pptx')
